#!/usr/bin/env python3
"""从KB1-KB8内容生成教学课件PPT - 干净简洁版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, re

# ==========================================
# 配色方案
# ==========================================
C = {
    "primary":   RGBColor(0x1A, 0x56, 0xDB),
    "secondary": RGBColor(0x37, 0x7A, 0xBF),
    "accent":    RGBColor(0xE8, 0x6A, 0x17),
    "bg_white":  RGBColor(0xFF, 0xFF, 0xFF),
    "bg_light":  RGBColor(0xF5, 0xF7, 0xFA),
    "text_dark": RGBColor(0x2D, 0x2D, 0x2D),
    "text_mid":  RGBColor(0x5A, 0x5A, 0x5A),
    "text_light":RGBColor(0x96, 0x96, 0x96),
    "green":     RGBColor(0x27, 0xAE, 0x60),
    "blue_bg":   RGBColor(0xDB, 0xEA, 0xFC),
    "orange_bg": RGBColor(0xFF, 0xE8, 0xD0),
    "yellow_bg": RGBColor(0xFF, 0xF3, 0xCD),
}

# ==========================================
# 辅助函数
# ==========================================
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=Pt(16), color=None, bold=False, align=PP_ALIGN.LEFT, name="Microsoft YaHei"):
    """添加文本框"""
    # Fix: if bold is not a bool (e.g. PP_ALIGN passed by mistake), swap
    if not isinstance(bold, bool):
        align = bold
        bold = False
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color or C["text_dark"]
        p.font.bold = bold
        p.font.name = name
        p.alignment = align
        if i > 0:
            p.space_before = Pt(4)
    return tf

def title_bar(slide, title, subtitle=""):
    """顶部标题栏"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = C["primary"]; bar.line.fill.background()
    tb(slide, 0.6, 0.15, 8.8, 0.55, title, Pt(28), C["bg_white"], True)
    if subtitle:
        tb(slide, 0.6, 0.7, 8.8, 0.35, subtitle, Pt(13), RGBColor(0xB0,0xC4,0xDE))

def kb_badge(slide, kb, left=0.3, top=0.15):
    """左上角KB标签"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.3), Inches(0.32))
    s.fill.solid(); s.fill.fore_color.rgb = C["bg_white"]; s.line.fill.background()
    p = s.text_frame.paragraphs[0]
    p.text = f"● {kb}"; p.font.size = Pt(10); p.font.color.rgb = C["primary"]; p.font.bold = True; p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER

def footer(slide, num):
    tb(slide, 0.3, 7.05, 5, 0.3, "AI时代能力培养 · 第1课：AI的发展与未来", Pt(8), C["text_light"])
    tb(slide, 8.5, 7.05, 1.2, 0.3, f"{num}/20", Pt(8), C["text_light"], align=PP_ALIGN.RIGHT)

def bullets(slide, items, left, top, width, height, size=Pt(14), spacing=Pt(14)):
    """项目符号列表"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item if item.startswith("•") else f"• {item}"
        p.font.size = size; p.font.color.rgb = C["text_dark"]; p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return tf

def card(slide, left, top, width, height, title, content_lines, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = C["bg_light"]
    s.line.color.rgb = border_color or C["secondary"]; s.line.width = Pt(1)
    tb(slide, left+0.15, top+0.08, width-0.3, 0.28, title, Pt(12), border_color or C["primary"], True)
    bullets(slide, content_lines, left+0.15, top+0.38, width-0.3, height-0.5, Pt(10), Pt(6))

# ==========================================
# 创建PPT
# ==========================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

pn = 0  # page number

# ==========================================
# 封面 (Slide 1)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.0))
rect.fill.solid(); rect.fill.fore_color.rgb = C["primary"]; rect.line.fill.background()
tb(s, 1, 0.5, 8, 0.8, "AI时代能力培养", Pt(42), C["bg_white"], True)
tb(s, 1, 1.3, 8, 0.6, "面向零基础学习者和企业管理者的AI入门课程", Pt(16), RGBColor(0xB0,0xC4,0xDE))
tb(s, 1, 2.2, 8, 0.5, "━━━━━━━━━━━━━━━━━━━━━━━━━━", Pt(12), RGBColor(0x80,0xA0,0xC0))
tb(s, 1, 3.6, 8, 0.7, "第一课：AI的发展与未来", Pt(34), C["primary"], True)
tb(s, 1, 4.4, 8, 0.5, "⏱ 总时长：约140分钟  |  8个知识块  |  理论+实操", Pt(14), C["text_mid"])
tb(s, 1, 5.2, 8, 0.9, "教学理念：会用AI → 懂AI → 会开发AI应用 → 了解AI前沿\n💼 企业视角 + 🎓 学生视角 双轨并行", Pt(13), C["text_mid"])
tb(s, 0.5, 6.9, 9, 0.3, "© 2026  |  讲义配套课件", Pt(8), C["text_light"], align=PP_ALIGN.CENTER)

# ==========================================
# 路线图 (Slide 2)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "本课路线图", "8个知识块 · 约140分钟")
kbs = [
    ("KB1", "课程导入\n与AI初识", "15min", "讲解+互动"),
    ("KB2", "AI发展简史", "25min", "故事化讲解"),
    ("KB3", "ChatGPT\n为什么火了", "15min", "讲解+讨论"),
    ("KB4", "大模型\n能做什么", "15min", "讲解+举例"),
    ("KB5", "AI会不会\n取代人", "15min", "互动讨论"),
    ("KB6", "主流AI\n工具介绍", "10min", "快速概览"),
    ("KB7", "动手实操", "40min", "注册+体验"),
    ("KB8", "总结与\n作业布置", "5min", "收尾"),
]
for i, (kb_id, title, dur, tag) in enumerate(kbs):
    left = 0.2 + i * 1.2
    s2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.6), Inches(1.1), Inches(2.3))
    s2.fill.solid(); s2.fill.fore_color.rgb = C["orange_bg"] if i == 6 else C["blue_bg"]
    s2.line.color.rgb = C["secondary"]; s2.line.width = Pt(1)
    tb(s, left+0.05, 1.68, 1.0, 0.22, kb_id, Pt(10), C["primary"], True, PP_ALIGN.CENTER)
    tb(s, left+0.05, 2.0, 1.0, 0.65, title, Pt(10), C["text_dark"], True, PP_ALIGN.CENTER)
    tb(s, left+0.05, 2.85, 1.0, 0.22, f"⏱ {dur}", Pt(9), C["accent"], False, PP_ALIGN.CENTER)
    tb(s, left+0.05, 3.15, 1.0, 0.35, tag, Pt(8), C["text_mid"], PP_ALIGN.CENTER)
    if i < 7:
        tb(s, left+1.05, 2.5, 0.2, 0.25, "→", Pt(14), C["text_light"], PP_ALIGN.CENTER)
tb(s, 0.5, 4.2, 9, 0.6, "⏱ 理论讲解：~95min  |  实操环节：~40min  |  🎓学生版+💼企业版 双轨并行", Pt(12), C["text_mid"], align=PP_ALIGN.CENTER)
footer(s, pn)

# ==========================================
# KB1 幻灯片 (Slide 3-4)
# ==========================================
# Slide 3: KB1 - AI是什么
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB1：课程导入与AI初识", "建立对AI的基本认知 · 15分钟")
kb_badge(s, "KB1 · 15min")
card(s, 0.4, 1.4, 4.4, 2.0, "🧠 AI是什么？",
    ["人工智能 = 让计算机像人一样 感知、思考、学习、决策",
     "核心类比：\"教小孩认猫\"",
     "→ 看很多猫的图片（数据）",
     "→ 告诉他\"这是猫\"（标签）",
     "→ 自己认出未见过品种（泛化）",
     "💼 企业视角：AI = 通用效率倍增杠杆",
     "不是行业，是渗透所有行业的基础能力"], C["primary"])
card(s, 5.2, 1.4, 4.4, 2.0, "🎯 课程学习理念",
    ["会用AI → 懂AI → 会开发AI应用 → 了解AI前沿",
     "理论30% + 实践70%",
     "\"不需成为程序员，要成为会用AI的人\"",
     "",
     "💼 企业管理者收获：",
     "一套给董事会的AI战略方案",
     "而非一个编程项目"], C["green"])
card(s, 0.4, 3.6, 9.2, 1.6, "🔗 AI > ML > DL > LLM（同心圆关系）",
    ["AI（人工智能）← 最大圈：所有让机器\"智能\"的技术",
     "  └─ ML（机器学习）← 从数据中学习规律",
     "      └─ DL（深度学习）← 多层神经网络",
     "          └─ LLM（大语言模型）← GPT/Claude/Gemini 等超大规模语言模型",
     "💼 企业学员：当供应商说\"基于最新LLM技术\"，追问\"具体基于哪个模型？和DeepSeek比有什么优势？\""], C["secondary"])
card(s, 0.4, 5.4, 9.2, 1.0, "⚡ 专用AI vs 通用AI",
    ["专用AI：下棋的不会翻译（AlphaGo/翻译模型/图像识别）",
     "通用AI（雏形）：一个模型能聊天/编程/翻译/写诗/做题（ChatGPT等）",
     "ChatGPT的意义：首次向大众展示\"通用人工智能\"的雏形"], C["accent"])
footer(s, pn)

# Slide 4: KB1 - 课程定位 + 互动
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB1：课程定位与互动", "")
kb_badge(s, "KB1 · 15min")
card(s, 0.4, 1.4, 4.4, 2.5, "🚲 电动自行车类比",
    ["AI是\"电动自行车\"，不是\"自动驾驶汽车\"",
     "🔋 它放大你的力量",
     "🧭 但方向的掌控始终在你手里",
     "  → 去哪、走哪条路、何时停 — 你决定",
     "",
     "💼 企业延伸：",
     "AI对企业是涡轮增压器，不是自动驾驶系统",
     "CEO仍需做战略决策，但支持性工作可加速3-10倍"], C["accent"])
card(s, 5.2, 1.4, 4.4, 2.5, "🙋 课堂互动",
    ["🎓 学生：\"有多少人用过ChatGPT/DeepSeek？\"",
     "🎓 学生：\"有多少人觉得AI会取代工作？\"",
     "",
     "💼 企业学员：\"过去一年，你的企业因AI做过",
     "    任何决策吗？采购/预算/战略讨论？\"",
     "",
     "核心理念：\"这门课给你的是判断力——",
     "学生判断AI对职业的影响，",
     "管理者判断AI对业务的影响。\"",
     "",
     "💬 企业讨论题：\"如果AI能把你的企业某一个",
     "    环节效率提升10倍，你希望是哪个环节？\""], C["primary"])
card(s, 0.4, 4.2, 9.2, 1.0, "📊 AI擅长 vs 不擅长",
    ["✅ 擅长：重复性任务 / 基于已有知识推理总结 / 快速生成多方案 / 24小时工作 / 多语言整合",
     "❌ 不擅长：情感共情 / 真正原创创造 / 物理实操经验 / 承担道德法律责任 / 理解隐性文化"], C["secondary"])
footer(s, pn)

# ==========================================
# KB2 幻灯片 (Slide 5-6)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB2：AI发展简史（1950-2017）", "一部跨越70年的\"冰与火之歌\" · 25分钟")
kb_badge(s, "KB2 · 25min")

timeline_left = [
    ("1950", "图灵测试 — \"机器能思考吗？\""),
    ("1956", "达特茅斯会议 — \"人工智能\"正式命名"),
    ("1973", "莱特希尔报告 → 第一次AI寒冬 ❄️"),
    ("1980s", "专家系统兴起 — \"如果…那么…\"规则"),
    ("1987", "专家系统崩溃 → 第二次AI寒冬 ❄️"),
]
for i, (year, event) in enumerate(timeline_left):
    y = 1.5 + i * 0.9
    tb(s, 0.4, y, 0.9, 0.25, year, Pt(14), C["primary"], True)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), Inches(y+0.3), Inches(0.03), Inches(0.45))
    line.fill.solid(); line.fill.fore_color.rgb = C["secondary"]; line.line.fill.background()
    tb(s, 1.4, y, 3.5, 0.25, event, Pt(11), C["text_dark"])

timeline_right = [
    ("1997", "深蓝击败卡斯帕罗夫 — AI首次超越人类智力竞技"),
    ("2012", "AlexNet夺冠 — 深度学习时代正式开启 🎆"),
    ("2016", "AlphaGo击败李世石 — 围棋\"最后堡垒\"被攻破"),
    ("2017", "⚡ Transformer论文发表 — \"Attention Is All You Need\""),
]
for i, (year, event) in enumerate(timeline_right):
    y = 1.5 + i * 1.1
    bg_c = C["yellow_bg"] if "2017" in year else C["bg_light"]
    s2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(y), Inches(4.5), Inches(0.9))
    s2.fill.solid(); s2.fill.fore_color.rgb = bg_c
    s2.line.color.rgb = C["accent"] if "2017" in year else C["secondary"]; s2.line.width = Pt(1)
    tb(s, 5.4, y+0.05, 4.1, 0.25, year, Pt(14), C["primary"], True)
    tb(s, 5.4, y+0.38, 4.1, 0.4, event, Pt(11), C["text_dark"])
tb(s, 0.5, 6.5, 9, 0.3, "⚠️ 2017年Transformer是改变一切的\"奇点\" — GPT/Claude/Gemini全部基于此架构", Pt(11), C["accent"], True)
footer(s, pn)

# Slide 6: KB2 - 2020s + 企业启示
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB2：大模型爆发 + 企业战略启示", "2020s-2026 · 商业视角")
kb_badge(s, "KB2 · 25min")
card(s, 0.4, 1.4, 4.4, 1.6, "🚀 2020s — 大模型爆发",
    ["2020 — GPT-3 (1750亿参数)",
     "2022.11 — ChatGPT发布 → 2个月破1亿用户",
     "2023 — GPT-4/Claude/Gemini → 百模大战",
     "2024 — 多模态成熟 + MCP协议发布",
     "2025 — Coding Agent成熟 (Claude Code/Codex)",
     "2026 — Agent工程化落地时代（现在）"], C["accent"])
card(s, 5.2, 1.4, 4.4, 1.6, "📌 四大关键启示",
    ["1️⃣ 72年技术积累，非一蹴而就",
     "2️⃣ 两次寒冬 = 过度乐观的代价",
     "3️⃣ Transformer = 现代AI的\"奇点\"",
     "4️⃣ 2022年 = AI民主化元年"], C["green"])
card(s, 0.4, 3.3, 9.2, 3.3, "💼 从AI发展史看企业战略：5条决策法则",
    ["启示1：警惕AI泡沫 — 历史上两次AI寒冬都因过度乐观。2026年同样存在泡沫风险。决策依据：ROI数据，而非行业噪音",
     "启示2：从专用AI到通用AI — 技术壁垒从\"模型能力\"转向\"数据+工程化+生态\"。1000万自建模型 vs 100万将通用AI与业务数据结合",
     "启示3：Transformer是奇点 — 底层技术趋同，差异化在应用层和数据。未来2-3年模型能力差异将缩小",
     "启示4：AI民主化 — 使用门槛消失，决策门槛上升。竞争优势从\"有没有AI\"转向\"如何更聪明地使用AI\"",
     "启示5：Agent工程化 — 2025-2026年Agent从概念走向落地。业务流程将被重构而非替代。先发优势窗口：6-18个月"], C["secondary"])
footer(s, pn)

# ==========================================
# KB3 幻灯片 (Slide 7)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB3：ChatGPT为什么火了？", "五大原因+五条企业决策法则 · 15分钟")
kb_badge(s, "KB3 · 15min")
reasons = [
    ("1️⃣ 能力泛化", "一个模型解决多种问题\n从\"专才\"到\"通才\"",
     "💼 通用平台吃专用工具\n警惕你买的AI方案被GPT-6覆盖"),
    ("2️⃣ 交互自然", "打字聊天 — 人人天生就会\n零学习成本 = 破圈关键",
     "💼 70%企业AI失败因员工不用\n用户推广成本≈零"),
    ("3️⃣ 门槛极低", "打开浏览器就能用\n免费+50+语言",
     "💼 定价从\"卖账号\"→\"卖成果\"\n按处理的文档数/节省工时收费"),
    ("4️⃣ 效果惊艳", "通过USMLE/BAR/编程竞赛\n大学水平论文写作",
     "💼 AI能力基准线快速上移\n3年规划以AI能力提升10倍为前提"),
    ("5️⃣ 时机成熟", "云计算+互联网数据+GPU算力\n疫情后数字化习惯+社交媒体传播",
     "💼 2026是黄金窗口期\n先发优势窗口：6-18个月"),
]
for i, (title, desc, biz) in enumerate(reasons):
    left = 0.15 + i * 1.95
    s2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(1.85), Inches(2.8))
    s2.fill.solid(); s2.fill.fore_color.rgb = C["bg_light"]; s2.line.color.rgb = C["secondary"]; s2.line.width = Pt(1)
    tb(s, left+0.1, 1.6, 1.65, 0.4, title, Pt(12), C["primary"], True)
    tb(s, left+0.1, 2.1, 1.65, 0.7, desc, Pt(9), C["text_dark"])
    tb(s, left+0.1, 2.95, 1.65, 0.9, biz, Pt(8), C["accent"])
card(s, 0.15, 4.55, 9.7, 0.7, "💡 ChatGPT最重要的贡献", ["让全人类（不仅是技术人员）都意识到AI时代的到来 —— \"AI的iPhone时刻\""], C["accent"])
card(s, 0.15, 5.45, 9.7, 0.7, "📊 增长数据", ["ChatGPT: 2个月破1亿用户  |  TikTok: 9个月  |  Instagram: 2.5年  →  史上增长最快的消费级应用"], C["secondary"])
footer(s, pn)

# ==========================================
# KB4 幻灯片 (Slide 8-9)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB4：大模型能够做什么？", "12项核心能力 + 个人应用场景 · 15分钟")
kb_badge(s, "KB4 · 15min")
caps = [
    ("📝 文本生成", "写文章/方案/报告/邮件"), ("💻 代码编写", "写代码/Debug/代码翻译"), ("🌍 翻译", "多语言互译·学术级"),
    ("📄 总结摘要", "长文提炼·300字总结20页"), ("🎓 知识问答", "私人导师·解释任何概念"), ("💡 创意生成", "头脑风暴·10个毕业选题"),
    ("📊 数据分析", "处理Excel·找异常·画图表"), ("🖼 图像理解", "识别图片·分析设计风格"), ("📽 PPT大纲", "生成演示文稿结构"),
    ("📚 文献综述", "梳理研究脉络·写综述"), ("📐 公式推导", "解释公式物理含义"), ("✓ 方案评估", "多维度分析方案优缺点"),
]
for i, (icon_title, desc) in enumerate(caps):
    col, row = i // 4, i % 4
    left, top = 0.3 + col * 3.2, 1.5 + row * 1.15
    tb(s, left, top, 3.0, 0.22, icon_title, Pt(11), C["primary"], True)
    tb(s, left, top+0.28, 3.0, 0.2, desc, Pt(9), C["text_mid"])
footer(s, pn)

# Slide 9: KB4 - 企业应用
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB4：企业级应用场景 + 分专业举例", "大模型能力对应企业ROI · 15分钟")
kb_badge(s, "KB4 · 15min")
card(s, 0.3, 1.4, 4.6, 3.2, "💼 企业职能 × AI能力矩阵",
    ["客服：7×24自动应答，成本降低40-60%",
     "营销：多平台文案批量生成，产出速度提升5-10倍",
     "研发：代码自动补全/检测/文档，效率提升30-50%",
     "法务：合同审查自动化，效率提升60-80%",
     "HR：简历筛选+面试问题生成，招聘效率提升40%",
     "财务：报表自动生成，月报3天→2小时",
     "高管办：行业分析/竞品追踪，调研效率提升5-8倍"], C["primary"])
card(s, 5.2, 1.4, 4.5, 3.2, "🎓 分专业应用举例",
    ["🔧 工程设计/机械",
     "设计初稿/技术文档/材料选择/实验数据/代码",
     "→ 竞标响应速度提升5倍",
     "",
     "🏗 土木工程",
     "施工方案/工程量估算/规范检索/项目报告",
     "→ 投标效率提升3-5倍",
     "",
     "🎨 建筑/艺术设计",
     "概念灵感/设计说明/Mood Board/色彩搭配",
     "→ 方案阶段2周缩短到2天"], C["secondary"])
card(s, 0.3, 4.8, 9.4, 1.4, "⚠️ 重要提醒",
    ["AI是\"辅助\"而非\"替代\"。AI生成内容必须审核，特别是涉及数据准确性、专业判断、安全的内容",
     "AI = \"初始草稿生成器\" + \"思路扩展器\" ≠ \"最终决策者\""], C["accent"])
footer(s, pn)

# ==========================================
# KB5 幻灯片 (Slide 10)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB5：AI会不会取代人？", "理性分析 + 企业劳动力战略 · 15分钟")
kb_badge(s, "KB5 · 15min")
card(s, 0.4, 1.4, 4.4, 2.0, "🎯 核心论点",
    ["\"AI取代的是具体任务，不是完整职业。\"",
     "\"会用AI的人（和企业）会取代不会用的。\"",
     "",
     "建筑设计师 = 100+项任务",
     "  AI能帮：画图/写方案/日照分析（5-10项）",
     "  AI做不到：客户沟通/审美判断/现场决策",
     "",
     "💼 对企业：AI使部分任务边际成本趋近于零",
     "McKinsey 2025: 60-70%工作被增强",
     "               15-20%被自动化",
     "                <5%完整职业3年内被取代"], C["accent"])
card(s, 5.2, 1.4, 4.4, 2.0, "🚲 电动自行车类比",
    ["AI是电动自行车，不是自动驾驶汽车",
     "🔋 它放大你的力量",
     "🧭 但方向的掌控始终在你手里",
     "",
     "💼 企业延伸：AI是涡轮增压器",
     "CEO仍需做战略决策",
     "但所有支持性工作可加速3-10倍",
     "",
     "正确态度：不是恐惧 → 而是学会骑它"], C["green"])
card(s, 0.4, 3.6, 9.2, 1.0, "💼 判断你企业中的\"高AI-Return任务\"（满足2/3即高ROI）",
    ["① 任务输入和输出可用文字/数字描述  |  ② 任务需要大量已有知识做判断  |  ③ 任务的错误可事后检查和修正"], C["secondary"])
card(s, 0.4, 4.8, 9.2, 1.8, "💼 给企业的四条战略建议",
    ["1. 做\"任务审计\"而非\"岗位审计\"——问\"这个岗位的哪些20%任务占了80%时间，且可用AI加速？\"",
     "2. \"人+AI\"组合的产出远超\"纯人\"或\"纯AI\"——组织设计的出发点是让人发挥判断力，AI处理规模化重复劳动",
     "3. 提前布局人才转型——培训现有员工使用AI工具，比外部招聘\"AI专家\"更具成本效益",
     "4. 建立AI使用规范和道德边界——明确哪些是\"AI辅助\"（AI建议+人决策），哪些是\"AI执行\"（AI自动+人审核）"], C["primary"])
footer(s, pn)

# ==========================================
# KB6 幻灯片 (Slide 11)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB6：主流AI工具介绍与选择", "六大工具 + 企业选型指南 · 10分钟")
kb_badge(s, "KB6 · 10min")
tools = [
    ("ChatGPT", "OpenAI", "综合最强\n生态完善"),
    ("Claude", "Anthropic", "编程最强\n长文分析"),
    ("Gemini", "Google", "多模态最强\n生态整合"),
    ("⭐DeepSeek", "深度求索", "完全免费\n中文极佳"),
    ("Kimi", "月之暗面", "200万字\n上下文"),
    ("Perplexity", "Perplexity AI", "AI搜索\n带引用"),
]
for i, (name, company, desc) in enumerate(tools):
    left = 0.12 + i * 1.62
    s2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(1.52), Inches(1.4))
    s2.fill.solid(); s2.fill.fore_color.rgb = C["yellow_bg"] if "DeepSeek" in name else C["bg_light"]
    s2.line.color.rgb = C["accent"] if "DeepSeek" in name else C["secondary"]; s2.line.width = Pt(2) if "DeepSeek" in name else Pt(1)
    tb(s, left+0.05, 1.55, 1.42, 0.25, name, Pt(11), C["primary"], True, PP_ALIGN.CENTER)
    tb(s, left+0.05, 1.8, 1.42, 0.2, company, Pt(7), C["text_light"], PP_ALIGN.CENTER)
    tb(s, left+0.05, 2.1, 1.42, 0.55, desc, Pt(9), C["text_dark"], PP_ALIGN.CENTER)
card(s, 0.12, 3.15, 4.8, 2.0, "🧭 个人选择决策树",
    ["写代码 → Claude / DeepSeek",
     "读论文/长文档 → Kimi（200万字）",
     "查最新信息 → Perplexity",
     "处理图片/视频 → Gemini",
     "国内方便 → DeepSeek（首选）",
     "全免费 → DeepSeek+Kimi+通义+豆包"], C["secondary"])
card(s, 5.2, 3.15, 4.7, 2.0, "💼 企业选型指南",
    ["全员日常 → DeepSeek（免费+直连）",
     "代码团队 → Claude Pro + DeepSeek",
     "文档处理 → Kimi（200万字上下文）",
     "数据安全 → 开源自部署模型",
     "实时信息 → Perplexity / Kimi联网",
     "⚠️ 核心原则：不要绑定单一供应商",
     "至少2-3个AI工具交叉验证"], C["accent"])
card(s, 0.12, 5.4, 9.7, 0.8, "🇨🇳 国内补充",
    ["通义千问（阿里）· 豆包（字节）· 文心一言（百度）— 全部免费·国内直连"], C["green"])
footer(s, pn)

# ==========================================
# KB7 幻灯片 (Slide 12-13)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB7：动手实操 — 注册与体验AI工具", "六步实操流程 + 学生/企业双版测试题 · 40分钟")
kb_badge(s, "KB7 · 40min")
steps = [
    ("Step 1", "注册ChatGPT\n15min"),
    ("Step 2", "注册Claude\n15min"),
    ("Step 3", "注册Gemini\n10min"),
    ("Step 4", "体验Perplexity\n10min"),
    ("Step 5", "⭐ 注册国内AI\n10min"),
    ("Step 6", "⭐ 对比测试\n最后完成"),
]
for i, (step, title) in enumerate(steps):
    left = 0.12 + i * 1.62
    s2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(1.52), Inches(1.3))
    s2.fill.solid(); s2.fill.fore_color.rgb = C["orange_bg"] if "⭐" in title else C["bg_light"]
    s2.line.color.rgb = C["accent"] if "⭐" in title else C["secondary"]; s2.line.width = Pt(2) if "⭐" in title else Pt(1)
    tb(s, left+0.05, 1.55, 1.42, 0.22, step, Pt(11), C["accent"], True, PP_ALIGN.CENTER)
    tb(s, left+0.05, 1.9, 1.42, 0.6, title, Pt(10), C["text_dark"], True, PP_ALIGN.CENTER)
    if i < 5:
        tb(s, left+1.38, 1.95, 0.25, 0.25, "→", Pt(12), C["text_light"], PP_ALIGN.CENTER)
card(s, 0.12, 3.1, 4.8, 2.0, "🔗 快速链接",
    ["ChatGPT: chatgpt.com",
     "Claude: claude.ai",
     "Gemini: gemini.google.com",
     "Perplexity: perplexity.ai",
     "⭐ DeepSeek: chat.deepseek.com",
     "Kimi: kimi.moonshot.cn",
     "通义千问: tongyi.aliyun.com",
     "豆包: doubao.com"], C["primary"])
card(s, 5.2, 3.1, 4.7, 2.0, "🎯 最低目标",
    ["✅ 每个学员至少注册1个AI工具",
     "   （推荐 DeepSeek — 免费+国内直连）",
     "✅ 完成对比测试：同一问题→至少2个AI",
     "✅ 记录对比结果",
     "",
     "⏱ 快 → 多注册几个工具深度体验",
     "⏱ 慢 → 只做 DeepSeek + 对比测试"], C["accent"])
card(s, 0.12, 5.3, 9.7, 0.8, "🆘 口诀",
    ["\"境外工具全跳过，DeepSeek必须搞\"  |  不要让任何学生在\"注册不了\"上消耗超过2分钟  |  立刻给替代方案"], C["accent"])
footer(s, pn)

# Slide 13: KB7 - 对比测试题
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB7：对比测试题（学生版+企业版）", "同一问题发给至少2个不同的AI · 40分钟")
kb_badge(s, "KB7 · 40min")
card(s, 0.3, 1.4, 4.6, 2.8, "🎓 学生版测试题",
    ["Q1：请用200字以内解释什么是\"结构力学\"。",
     "",
     "Q2：一个工程设计师在日常工作中",
     "    如何使用AI提高效率？列出5个具体场景。",
     "",
     "Q3：用Python写一个计算BMI的小程序",
     "    （输入身高m和体重kg，",
     "      输出BMI值和健康分级）。"], C["primary"])
card(s, 5.2, 1.4, 4.5, 2.8, "💼 企业版测试题",
    ["Q1：请用200字以内向企业管理者解释",
     "    什么是\"大语言模型（LLM）\"，",
     "    重点说明它能解决什么商业问题。",
     "",
     "Q2：一家50人的科技公司应该如何",
     "    在日常运营中使用AI提高效率？",
     "    列出5个最具ROI的场景和预期效果。",
     "",
     "Q3：请起草一份200字的内部邮件，",
     "    向团队宣布公司将开始使用AI工具"],
        C["accent"])
card(s, 0.3, 4.5, 9.4, 1.8, "📋 对比记录维度",
    ["回答质量(1-5)  |  回答实用性(1-5)  |  代码/输出正确性  |  回答速度  |  中文自然度  |  使用便捷度  |  长期使用理由",
     "💼 企业学员额外：\"哪个工具的回答最适合直接用于内部沟通？哪个最需要修改？\""], C["secondary"])
footer(s, pn)

# ==========================================
# KB8 幻灯片 (Slide 14-15)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB8：课程总结 · 课后作业 · 下次课预告", "5分钟收尾")
kb_badge(s, "KB8 · 5min")

card(s, 0.3, 1.4, 9.4, 1.6, "📌 本课五句话核心回顾",
    ["1️⃣ AI发展是一条70年的路，不是一夜之间冒出来的",
     "2️⃣ ChatGPT = 技术突破 + 产品设计 + 时代机遇的完美结合",
     "3️⃣ 大模型能做的事情远超你想象 — 但需要你学会\"指挥\"它",
     "4️⃣ AI是电动自行车，不是自动驾驶 — 你掌握方向",
     "5️⃣ 不要只用一种AI — 组建你的\"AI顾问团\""], C["primary"])

card(s, 0.3, 3.2, 4.5, 1.8, "🎓 学生版作业",
    ["作业1（必做）：AI工具对比笔记",
     "  3个问题×3个工具，≥400字",
     "作业2（必做）：传统vs AI效率对比",
     "  选本专业真实问题，两种方式解决",
     "作业3（选做）：AI应用场景头脑风暴",
     "  未来工作中5+个AI提效场景"],
        C["primary"])
card(s, 5.2, 3.2, 4.5, 1.8, "💼 企业版作业",
    ["作业1（必做）：AI工具对比笔记",
     "  使用企业版测试问题",
     "作业2（必做）：企业任务审计",
     "  10个核心任务×AI替代程度→1页矩阵",
     "作业3（选做）：AI供应商初筛报告",
     "  3个候选供应商×5维评估"],
        C["accent"])
card(s, 0.3, 5.2, 9.4, 0.8, "🔮 下次课预告",
    ["Prompt Engineering（提示词工程）— RCTE框架 · 论文写作 · 代码编写 · PPT制作 · 方案设计 · 全部AI辅助"],
        C["accent"])
footer(s, pn)

# Slide 15: KB8 - 关键概念速查
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "附录：关键概念速查表", "讲课过程中如有学生提问概念，可快速翻到此页")
concepts = [
    ("AI（人工智能）", "让计算机像人一样感知/思考/学习/决策", "通用效率杠杆，非独立行业"),
    ("图灵测试", "对话中分不清人还是机器", "AI交互已通过图灵测试"),
    ("Transformer", "2017年架构，所有现代LLM的基础", "底层技术趋同，差异化在上层"),
    ("大语言模型(LLM)", "海量文本训练的超级神经网络", "企业AI的\"通用引擎\""),
    ("ChatGPT", "OpenAI的对话AI，2022年引爆全球", "AI的iPhone时刻"),
    ("Prompt(提示词)", "你发给AI的指令/问题", "精准传达需求=企业核心竞争力"),
    ("多模态", "AI同时处理文本/图片/音频/视频", "可处理合同扫描件/会议录音"),
    ("AI幻觉", "AI一本正经编造假信息", "关键决策必须有审核机制"),
    ("Agent(智能体)", "能自主规划、使用工具的AI", "未来的\"数字员工\""),
    ("MCP", "Model Context Protocol", "AI时代的USB-C接口标准"),
]
for i, (concept, desc_stu, desc_biz) in enumerate(concepts):
    col = i // 5; row = i % 5
    left = 0.2 + col * 5.0; top = 1.5 + row * 0.95
    tb(s, left, top, 4.7, 0.22, f"■ {concept}", Pt(11), C["primary"], True)
    tb(s, left, top+0.25, 2.2, 0.2, f"🎓 {desc_stu}", Pt(8), C["text_dark"])
    tb(s, left+2.3, top+0.25, 2.4, 0.2, f"💼 {desc_biz}", Pt(8), C["accent"])
footer(s, pn)

# ==========================================
# KB2企业启示补充 (Slide 16)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB2·KB3·KB5：企业核心知识点汇总", "AI发展史商业启示 + ChatGPT商业法则 + 劳动力战略")
card(s, 0.2, 1.3, 3.1, 5.2, "📜 KB2: AI发展史5条启示",
    ["1. 警惕泡沫 — ROI决策",
     "2. 通用AI — 采购策略转变",
     "3. Transformer — 底层趋同",
     "4. AI民主化 — 决策门槛↑",
     "5. Agent工程化 — 窗口6-18月"], C["primary"])
card(s, 3.45, 1.3, 3.1, 5.2, "💡 KB3: ChatGPT五条商业法则",
    ["1. 通用平台吃专用工具",
     "2. 70%失败因员工不用",
     "3. 定价从账号→成果",
     "4. 规划以10倍提升为前提",
     "5. 2026黄金窗口期"], C["accent"])
card(s, 6.7, 1.3, 3.1, 5.2, "🏢 KB5: 企业劳动力战略",
    ["任务审计 vs 岗位审计",
     "人+AI > 纯人或纯AI",
     "培训现有员工 > 外招",
     "建立AI使用规范",
     "",
     "高ROI任务:",
     "文字描述+知识判断",
     "+错误可事后修正"], C["green"])
footer(s, pn)

# ==========================================
# KB4企业应用补充 (Slide 17)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
title_bar(s, "KB4·KB6·KB7：工具与企业落地要点", "能力矩阵 + 选型指南 + 实操要点")
card(s, 0.2, 1.3, 4.7, 3.0, "💼 KB4: 企业级场景ROI速查",
    ["客服：7×24自动，成本↓40-60%",
     "营销：内容产出速度↑5-10倍",
     "研发：开发效率↑30-50%",
     "法务：合同审查效率↑60-80%",
     "HR：招聘效率↑40%",
     "财务：月报3天→2小时",
     "高管：调研效率↑5-8倍",
     "",
     "⚠️ AI = 辅助 ≠ 替代，关键决策需人工审核"], C["primary"])
card(s, 5.2, 1.3, 4.6, 3.0, "💼 KB6: 企业AI工具选型矩阵",
    ["全员日常→DeepSeek（免费+直连）",
     "代码团队→Claude Pro+DeepSeek",
     "文档处理→Kimi（200万字）",
     "数据安全→开源自部署",
     "实时信息→Perplexity/Kimi联网",
     "",
     "⚠️ 核心原则：",
     "不要绑定单一供应商",
     "至少2-3个AI交叉验证"], C["accent"])
card(s, 0.2, 4.5, 9.6, 1.4, "💼 KB7: 企业版实操要点",
    ["企业学员测试Prompt：①向管理者解释LLM的商业价值  ②50人公司AI提效场景  ③起草团队AI工具采用邮件",
     "教师巡查时确认：企业学员是否在使用企业版测试题？Prompt是否足够具体？",
     "对比测试记录额外维度：哪个工具的回答最适合直接用于内部沟通？哪个最需要修改？"], C["secondary"])
footer(s, pn)

# ==========================================
# 结束页 (Slide 18)
# ==========================================
pn += 1; s = blank_slide(prs); set_bg(s, C["bg_white"])
rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.2), Inches(10), Inches(2.3))
rect.fill.solid(); rect.fill.fore_color.rgb = C["primary"]; rect.line.fill.background()
tb(s, 1, 1.5, 8, 1.0, "72年前，图灵问：机器能思考吗？\n今天，你自己来回答这个问题。", Pt(22), C["text_dark"], align=PP_ALIGN.CENTER)
tb(s, 1, 2.8, 8, 0.8, "欢迎大家来到AI时代\n这一次，你不是旁观者。", Pt(26), C["primary"], True, align=PP_ALIGN.CENTER)
tb(s, 1, 5.5, 8, 0.5, "🎓 对学生：你不是在学编程，你是在学未来最重要的生存技能", Pt(15), C["bg_white"], align=PP_ALIGN.CENTER)
tb(s, 1, 6.0, 8, 0.5, "💼 对企业学员：你不是在选要不要AI，而是在选主动驾驭它还是被动应对它", Pt(15), C["bg_white"], align=PP_ALIGN.CENTER)
tb(s, 1, 6.7, 8, 0.6, "📅 下次课：Prompt Engineering  |  📋 作业截止：下次上课前一天22:00", Pt(13), RGBColor(0xC0,0xD0,0xE0), align=PP_ALIGN.CENTER)

# ==========================================
# 保存
# ==========================================
out = "d:/2026/AI-Learning-Route/讲义/第1周/AI的发展与未来/AI的发展与未来-教学课件_v2.pptx"
prs.save(out)
print(f"DONE: {os.path.getsize(out)/1024:.1f}KB, {len(prs.slides)} slides")
