#!/usr/bin/env python3
"""生成第1课教学辅助PPT — 教学课件风格"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==========================================
# 配色方案 — 教学课件风（温暖、清晰、护眼）
# ==========================================
COLORS = {
    "primary":    RGBColor(0x1A, 0x56, 0xDB),   # 深蓝 — 标题
    "secondary":  RGBColor(0x37, 0x7A, 0xBF),   # 中蓝 — 副标题
    "accent":     RGBColor(0xE8, 0x6A, 0x17),   # 橙色 — 强调/互动提示
    "bg_light":   RGBColor(0xF5, 0xF7, 0xFA),   # 浅灰蓝 — 背景
    "bg_white":   RGBColor(0xFF, 0xFF, 0xFF),   # 白色
    "text_dark":  RGBColor(0x2D, 0x2D, 0x2D),   # 深灰 — 正文
    "text_mid":   RGBColor(0x5A, 0x5A, 0x5A),   # 中灰 — 辅助文字
    "text_light": RGBColor(0x96, 0x96, 0x96),   # 浅灰
    "green":      RGBColor(0x27, 0xAE, 0x60),   # 绿色 — 正面/完成
    "red":        RGBColor(0xE7, 0x4C, 0x3C),   # 红色 — 强调/警告
    "yellow_bg":  RGBColor(0xFF, 0xF3, 0xCD),   # 淡黄 — 高亮框
    "blue_bg":    RGBColor(0xDB, 0xEA, 0xFC),   # 淡蓝 — 信息框
    "orange_bg":  RGBColor(0xFF, 0xE8, 0xD0),   # 淡橙 — 提示框
}

# ==========================================
# 辅助函数
# ==========================================
def set_slide_bg(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(18),
                color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color or COLORS["text_dark"]
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_kb_badge(slide, kb_label, duration, left=0.3, top=0.15):
    """在左上角添加KB标签"""
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(1.6), Inches(0.35)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = COLORS["primary"]
    tag.line.fill.background()
    tf = tag.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f" {kb_label} · {duration}"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["bg_white"]
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

def add_teacher_tip(slide, text, left=0.3, top=6.6, width=9.4, height=0.45):
    """添加教师提示条（淡橙色底）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["orange_bg"]
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"💡 教师提示：{text}"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["accent"]
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT

def add_footer(slide, slide_num, total=16):
    """添加页脚（页码+课程名）"""
    add_textbox(slide, 0.3, 7.0, 5, 0.3,
                f"AI时代能力培养 · 第1课：AI的发展与未来",
                Pt(8), COLORS["text_light"])
    add_textbox(slide, 8.5, 7.0, 1.5, 0.3,
                f"{slide_num}/{total}",
                Pt(8), COLORS["text_light"], alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle="", kb_label="", duration=""):
    """添加统一的节标题"""
    # 顶部蓝色横条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()

    # 标题
    add_textbox(slide, 0.5, 0.3, 9, 0.7, title,
                Pt(32), COLORS["primary"], bold=True)
    if subtitle:
        add_textbox(slide, 0.5, 0.9, 9, 0.5, subtitle,
                    Pt(14), COLORS["text_mid"])
    if kb_label:
        add_kb_badge(slide, kb_label, duration)

def add_bullet_list(slide, items, left=0.5, top=1.8, width=8.5, font_size=Pt(14),
                    color=None, spacing=Pt(24)):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5 - top))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = font_size
        p.font.color.rgb = color or COLORS["text_dark"]
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return tf

def add_card(slide, left, top, width, height, title, content, color=None):
    """添加卡片样式的内容块"""
    # 卡片背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["bg_light"]
    shape.line.color.rgb = color or COLORS["secondary"]
    shape.line.width = Pt(1)
    # 标题
    add_textbox(slide, left + 0.15, top + 0.08, width - 0.3, 0.3,
                title, Pt(13), color or COLORS["primary"], bold=True)
    # 内容
    add_textbox(slide, left + 0.15, top + 0.38, width - 0.3, height - 0.5,
                content, Pt(10), COLORS["text_dark"])

# ==========================================
# 创建演示文稿
# ==========================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ==========================================
# Slide 1: 封面
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, COLORS["bg_white"])

# 大色块装饰
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = COLORS["primary"]
shape.line.fill.background()

add_textbox(slide, 1, 0.6, 8, 0.8, "AI时代能力培养",
            Pt(42), COLORS["bg_white"], bold=True)
add_textbox(slide, 1, 1.5, 8, 0.6, "面向零基础大学生的AI入门速成课程",
            Pt(18), RGBColor(0xB0, 0xC4, 0xDE))
add_textbox(slide, 1, 2.3, 8, 0.5, "━━━━━━━━━━━━━━━━━━━━━━━━━━",
            Pt(12), RGBColor(0x80, 0xA0, 0xC0))

add_textbox(slide, 1, 3.8, 8, 0.8, "第一课：AI的发展与未来",
            Pt(36), COLORS["primary"], bold=True)
add_textbox(slide, 1, 4.6, 8, 0.5, "⏱ 总时长：约140分钟  |  8个知识块  |  理论讲解+动手实操",
            Pt(14), COLORS["text_mid"])

# 底部信息
add_textbox(slide, 1, 5.6, 8, 1.0,
            "教学理念：会用AI → 懂AI → 会开发AI应用 → 了解AI前沿\n"
            "📍 课前准备：浏览器 + 邮箱 + DeepSeek账号",
            Pt(13), COLORS["text_mid"])

add_textbox(slide, 0.5, 7.0, 9.5, 0.3,
            "第1课 · AI的发展与未来  |  V1.0  |  © 2026",
            Pt(8), COLORS["text_light"], alignment=PP_ALIGN.CENTER)

# ==========================================
# Slide 2: 本课路线图
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "本课路线图", "8个知识块 · 约140分钟", "总览", "140min")

# 路线图 — 用8个卡片排列
kb_items = [
    ("KB1", "课程导入\n与AI初识", "15min", "讲解+互动"),
    ("KB2", "AI发展简史", "25min", "故事化讲解"),
    ("KB3", "ChatGPT\n为什么火了", "15min", "讲解+讨论"),
    ("KB4", "大模型\n能做什么", "15min", "讲解+举例"),
    ("KB5", "AI会不会\n取代人", "15min", "互动讨论"),
    ("KB6", "主流AI\n工具介绍", "10min", "快速概览"),
    ("KB7", "动手实操", "40min", "注册+体验"),
    ("KB8", "总结与\n作业布置", "5min", "收尾"),
]

for i, (kb_id, title, duration, tag) in enumerate(kb_items):
    left = 0.3 + i * 1.18
    top = 1.8
    # 卡片
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(1.08), Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["blue_bg"] if i != 6 else COLORS["orange_bg"]
    shape.line.color.rgb = COLORS["secondary"]
    shape.line.width = Pt(1)
    # KB标签
    add_textbox(slide, left + 0.05, top + 0.08, 1.0, 0.25,
                kb_id, Pt(10), COLORS["primary"], bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_textbox(slide, left + 0.05, top + 0.45, 1.0, 0.7,
                title, Pt(11), COLORS["text_dark"], bold=True, alignment=PP_ALIGN.CENTER)
    # 时长
    add_textbox(slide, left + 0.05, top + 1.3, 1.0, 0.25,
                f"⏱ {duration}", Pt(9), COLORS["accent"], alignment=PP_ALIGN.CENTER)
    # 类型
    add_textbox(slide, left + 0.05, top + 1.65, 1.0, 0.4,
                tag, Pt(8), COLORS["text_mid"], alignment=PP_ALIGN.CENTER)
    # 箭头
    if i < 7:
        add_textbox(slide, left + 0.95, top + 0.85, 0.28, 0.3,
                    "→", Pt(16), COLORS["text_light"], alignment=PP_ALIGN.CENTER)

# 底部时间线
add_textbox(slide, 0.5, 4.3, 9, 0.4,
            "⏱ 理论讲解：~95min  |  实操环节：~40min  |  机动时间：~5min",
            Pt(12), COLORS["text_mid"], alignment=PP_ALIGN.CENTER)

add_teacher_tip(slide, "这张PPT在讲课时保持打开，提醒学生和您当前的进度位置。每讲完一个KB，翻回此页指出已完成的部分。")
add_footer(slide, 2)

# ==========================================
# Slide 3: KB1 要点 (上)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB1：课程导入与AI初识", "建立对AI的基本认知 + 激发学习兴趣", "KB1", "15min")
add_kb_badge(slide, "KB1", "15min")

# AI定义卡片
add_card(slide, 0.4, 1.5, 4.4, 1.4,
         "🧠 AI是什么？",
         "人工智能 = 让计算机像人一样\n\n"
         "👁 感知 → 🧠 思考 → 📚 学习 → ✅ 决策\n\n"
         "💬 核心类比：\"教小孩认猫\"",
         COLORS["primary"])

add_card(slide, 5.2, 1.5, 4.4, 1.4,
         "🎯 课程学习理念",
         "✦ 会用AI → 懂AI → 会开发AI应用 → 了解AI前沿\n"
         "✦ 理论30% + 实践70%\n"
         "✦ \"你不是学编程，你是学会用AI解决专业问题\"",
         COLORS["green"])

# AI/ML/DL/LLM关系
add_card(slide, 0.4, 3.2, 9.2, 1.6,
         "🔗 AI > ML > DL > LLM （同心圆关系）",
         "┌─────────────────────────────────────────┐\n"
         "│  AI 人工智能  ← 最大的圈：所有让机器\"智能\"的技术      │\n"
         "│  ┌──────────────────────────────┐        │\n"
         "│  │  ML 机器学习  ← 从数据中学习规律          │        │\n"
         "│  │  ┌────────────────────┐      │        │\n"
         "│  │  │  DL 深度学习  ← 多层神经网络    │      │        │\n"
         "│  │  │  ┌──────────┐     │      │        │\n"
         "│  │  │  │ LLM 大模型│  ← 超大规模语言模型 │        │\n"
         "│  │  │  │ GPT/Claude│     │      │        │\n"
         "│  │  │  └──────────┘     │      │        │\n"
         "│  │  └────────────────────┘      │        │\n"
         "│  └──────────────────────────────┘        │\n"
         "└─────────────────────────────────────────┘\n"
         "📌 记忆法：AI是大超市，ML是食品区，DL是生鲜区，LLM是和牛专柜",
         COLORS["secondary"])

add_teacher_tip(slide, "此处可提问：\"你用过哪些AI工具？\" 快速举手调查，活跃气氛。预计5分钟互动。")
add_footer(slide, 3)

# ==========================================
# Slide 4: KB1 要点 (下) + KB2 过渡
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB1 · AI与人类的对比 + 课程定位", "", "KB1", "15min")

add_card(slide, 0.4, 1.5, 4.4, 2.0,
         "⚡ AI擅长 vs 不擅长",
         "✅ 擅长：\n"
         "  • 大量重复性任务（永不疲倦）\n"
         "  • 基于已有知识推理总结\n"
         "  • 快速生成多种方案变体\n"
         "  • 多语言、多领域知识整合\n\n"
         "❌ 不擅长：\n"
         "  • 真正原创性的创造\n"
         "  • 需要物理世界经验的操作\n"
         "  • 承担道德责任\n"
         "  • 理解组织的隐性文化",
         COLORS["accent"])

add_card(slide, 5.2, 1.5, 4.4, 2.0,
         "🚲 核心类比：电动自行车",
         "\"AI是电动自行车，\n  不是自动驾驶汽车\"\n\n"
         "🔋 它放大你的力量\n"
         "🧭 但方向始终由你掌控\n\n"
         "正确的态度：\n"
         "不是\"恐惧AI取代我\"\n"
         "而是\"学会骑这辆电动自行车\"",
         COLORS["green"])

add_card(slide, 0.4, 3.8, 9.2, 1.0,
         "📊 专用AI vs 通用AI",
         "专用AI：下棋的不会翻译 → AlphaGo / 翻译模型 / 图像识别模型\n"
         "通用AI（雏形）：一个模型解决多种任务 → ChatGPT 能聊天/编程/翻译/写诗/做题……\n"
         "ChatGPT的意义：首次向大众展示了\"通用人工智能\"的雏形",
         COLORS["secondary"])

add_teacher_tip(slide, "过渡到KB2：\"那么AI是怎么从专用走到今天的？让我们回到70年前……\" 翻下一页。")
add_footer(slide, 4)

# ==========================================
# Slide 5: KB2 AI发展简史时间线 (上)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB2：AI发展简史（1950-2026）", "一部跨越70年的\"冰与火之歌\"", "KB2", "25min")

# 时间线左列
timeline_left = [
    ("1950", "图灵测试", "\"机器能思考吗？\""),
    ("1956", "达特茅斯会议", "\"人工智能\"正式命名"),
    ("1973", "莱特希尔报告", "→ 第一次AI寒冬 ❄️"),
    ("1980s", "专家系统兴起", "\"如果…那么…\"规则"),
    ("1987", "专家系统崩溃", "→ 第二次AI寒冬 ❄️"),
]

for i, (year, title, desc) in enumerate(timeline_left):
    y = 1.6 + i * 0.95
    # 年份标记
    add_textbox(slide, 0.4, y, 0.8, 0.3, year,
                Pt(14), COLORS["primary"], bold=True)
    # 竖线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.15), Inches(y + 0.35), Inches(0.03), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["secondary"]
    shape.line.fill.background()
    # 事件
    add_textbox(slide, 1.4, y, 3.5, 0.3, title,
                Pt(12), COLORS["text_dark"], bold=True)
    add_textbox(slide, 1.4, y + 0.3, 3.5, 0.3, desc,
                Pt(10), COLORS["text_mid"])

# 时间线右列
timeline_right = [
    ("1997", "深蓝击败卡斯帕罗夫", "AI首次在智力竞技上超越人类"),
    ("2012", "AlexNet夺冠", "🎆 深度学习时代正式开启"),
    ("2016", "AlphaGo击败李世石", "围棋\"最后的人类堡垒\"被攻破"),
    ("2017", "⚡ Transformer论文发表", "\"Attention Is All You Need\" — 改变AI历史的8页论文"),
]

for i, (year, title, desc) in enumerate(timeline_right):
    y = 1.6 + i * 1.15
    bg_color = COLORS["yellow_bg"] if "2017" in year else COLORS["bg_light"]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(y), Inches(4.5), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = COLORS["accent"] if "2017" in year else COLORS["secondary"]
    shape.line.width = Pt(1)
    add_textbox(slide, 5.4, y + 0.05, 1.2, 0.25, year,
                Pt(14), COLORS["primary"], bold=True)
    add_textbox(slide, 6.3, y + 0.05, 3.2, 0.25, title,
                Pt(11), COLORS["text_dark"], bold=True)
    add_textbox(slide, 5.4, y + 0.45, 4.1, 0.4, desc,
                Pt(9), COLORS["text_mid"])

add_teacher_tip(slide, "KB2核心教学策略：不是念时间线——而是讲\"故事\"。重点强调Transformer（2017）是改变一切的节点。")
add_footer(slide, 5)

# ==========================================
# Slide 6: KB2 AI发展简史 (下) 2020s + 关键启示
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB2：大模型爆发（2020s-2026）+ 关键启示", "", "KB2", "25min")

add_card(slide, 0.4, 1.5, 4.4, 1.6,
         "🚀 2020s — 大模型爆发",
         "2020 — GPT-3 (1750亿参数)\n"
         "2022.11 — ChatGPT发布 🔥\n"
         "          → 2个月破1亿用户\n"
         "2023 — GPT-4/Claude/Gemini → 百模大战\n"
         "2024 — 多模态成熟 + MCP协议\n"
         "2025 — Coding Agent成熟\n"
         "         Claude Code/Codex\n"
         "2026 — Agent工程化落地时代",
         COLORS["accent"])

add_card(slide, 5.2, 1.5, 4.4, 1.6,
         "📌 四大关键启示",
         "1️⃣ AI不是一蹴而就 — 72年积累\n"
         "2️⃣ 两次寒冬教会我们谦逊\n"
         "   过度乐观 → 现实打击 → 经费削减\n"
         "3️⃣ Transformer是\"奇点\"\n"
         "   2017年的8页论文开启现代LLM时代\n"
         "4️⃣ 2022 = AI民主化元年\n"
         "   ChatGPT之前：AI属于少数精英\n"
         "   ChatGPT之后：AI属于所有人",
         COLORS["green"])

add_card(slide, 0.4, 3.4, 9.2, 1.6,
         "🎯 1950-2026 AI发展时间线全景速查",
         "1950       1956     1973    1980s    1987     1997        2012      2016      2017         2020      2022.11      2026\n"
         "图灵测试 → 达特茅斯 → 寒冬1 → 专家系统 → 寒冬2 → 深蓝夺冠 → AlexNet → AlphaGo → ⚡Transformer → GPT-3 → ChatGPT 🔥 → Agent时代\n"
         "  AI诞生    命名      ❄️    小春天     ❄️    专用AI      DL崛起    围棋    改变一切      大模型     AI民主化      工程落地",
         COLORS["secondary"])

add_teacher_tip(slide, "此处可做互动：\"请把左边的事件和右边的年份连线\"（PPT翻页后展示答案）。预计2分钟。")
add_footer(slide, 6)

# ==========================================
# Slide 7: KB3 ChatGPT为什么火了
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB3：ChatGPT为什么火了？", "五大原因 —— 一场\"完美风暴\"", "KB3", "15min")

reasons = [
    ("1️⃣ 能力泛化", "一个模型解决多种问题\n从\"专才\"到\"通才\"", "瑞士军刀 → 全能助手"),
    ("2️⃣ 交互自然", "打字聊天 — 每个人天生就会\n零学习成本 = 破圈关键", "不需要学任何命令/API"),
    ("3️⃣ 门槛极低", "打开浏览器就能用\n免费版可用 支持50+语言", "对比：以前AI需GPU+编程"),
    ("4️⃣ 效果惊艳", "通过USMLE医学执照\n通过BAR律师资格\n多项编程竞赛排名前列", "\"AI的智能到底有多高？\""),
    ("5️⃣ 时机成熟", "云计算普及 + 互联网数据积累\nGPU算力1000倍提升\n社交媒体病毒传播", "技术+产品+时机的完美结合"),
]

for i, (title, desc, tag) in enumerate(reasons):
    left = 0.2 + i * 1.92
    top = 1.7
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(1.8), Inches(2.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["bg_light"]
    shape.line.color.rgb = COLORS["secondary"]
    shape.line.width = Pt(1)
    add_textbox(slide, left + 0.1, top + 0.1, 1.6, 0.5, title,
                Pt(13), COLORS["primary"], bold=True)
    add_textbox(slide, left + 0.1, top + 0.65, 1.6, 1.0, desc,
                Pt(9), COLORS["text_dark"])
    add_textbox(slide, left + 0.1, top + 1.85, 1.6, 0.35, f"→ {tag}",
                Pt(8), COLORS["text_mid"])

# 底部关键洞察
add_card(slide, 0.4, 4.3, 9.2, 0.8,
         "💡 关键洞察",
         "ChatGPT最重要的贡献：让全人类（不仅是技术人员）都意识到AI时代的到来 —— \"AI的iPhone时刻\"",
         COLORS["accent"])

add_card(slide, 0.4, 5.3, 9.2, 0.8,
         "📊 增长对比",
         "ChatGPT: 2个月破1亿用户  |  TikTok: 9个月  |  Instagram: 2.5年  |  →  史上增长最快的消费级应用",
         COLORS["secondary"])

add_teacher_tip(slide, "此处可互动：\"过去一年，ChatGPT（或类似AI）改变了你的生活或学习吗？\" 1分钟快速分享。")
add_footer(slide, 7)

# ==========================================
# Slide 8: KB4 大模型能力全景
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB4：大模型能够做什么？", "12项核心能力 + 你的专业场景", "KB4", "15min")

capabilities = [
    ("📝 文本生成", "写文章/方案/报告/邮件"),
    ("💻 代码编写", "写代码/Debug/代码翻译"),
    ("🌍 翻译", "多语言互译·学术级"),
    ("📄 总结摘要", "长文提炼要点·300字总结20页"),
    ("🎓 知识问答", "私人导师·解释任何概念"),
    ("💡 创意生成", "头脑风暴·10个毕业选题"),
    ("📊 数据分析", "处理Excel·找异常·画图表"),
    ("🖼 图像理解", "识别图片·分析设计风格"),
    ("📽 PPT大纲", "生成演示文稿结构"),
    ("📚 文献综述", "梳理研究脉络·写综述"),
    ("📐 公式推导", "解释公式的物理含义"),
    ("✓ 方案评估", "多维度分析方案优缺点"),
]

for i, (icon_title, desc) in enumerate(capabilities):
    col = i // 4
    row = i % 4
    left = 0.4 + col * 3.2
    top = 1.5 + row * 1.2
    add_textbox(slide, left, top, 3.0, 0.25,
                icon_title, Pt(12), COLORS["primary"], bold=True)
    add_textbox(slide, left, top + 0.3, 3.0, 0.2,
                desc, Pt(9), COLORS["text_mid"])

add_teacher_tip(slide, "此处应放慢节奏，挑3-4个与班上专业最相关的能力展开举例。12个能力不用全念。")
add_footer(slide, 8)

# ==========================================
# Slide 9: KB4 分专业应用场景
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB4：分专业AI应用场景举例", "让学生从自己专业的角度看到AI的实用性", "KB4", "15min")

add_card(slide, 0.4, 1.5, 2.9, 4.0,
         "🔧 工程设计/机械",
         "• 课程设计初始方案构思\n• 技术文档撰写和润色\n• 材料选择初步建议\n• 实验数据整理和分析\n• 代码编写（Python/MATLAB）\n\n💬 \"帮我写一份齿轮减速器\n  设计计算说明书\"",
         COLORS["primary"])

add_card(slide, 3.55, 1.5, 2.9, 4.0,
         "🏗 土木工程",
         "• 施工组织设计方案框架\n• 工程量辅助估算\n• 规范条文快速检索\n• 项目报告模板生成\n• 结构分析概念辅助学习\n\n💬 \"解释一下这个结构力学\n  公式每一步的物理含义\"",
         COLORS["secondary"])

add_card(slide, 6.7, 1.5, 2.9, 4.0,
         "🎨 建筑/艺术设计",
         "• 设计概念灵感发散\n• 设计说明文档撰写\n• Mood Board色彩搭配\n• 设计风格分析\n• 作品集文案优化\n\n💬 \"给我5个现代简约风格的\n  咖啡厅设计概念方向\"",
         COLORS["accent"])

add_teacher_tip(slide, "这是KB4最关键的部分 — 学生看到自己专业的名字会自然集中注意力。根据班级实际专业分布调整展开深度。")
add_footer(slide, 9)

# ==========================================
# Slide 10: KB5 AI会不会取代人
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB5：AI会不会取代人？", "坦诚面对焦虑，理性分析 + 行动方向", "KB5", "15min")

add_card(slide, 0.4, 1.5, 4.4, 1.8,
         "🎯 核心论点",
         "\"AI取代的是具体任务，\n  不是完整职业。\"\n\n"
         "例：建筑设计师 = 100+项任务\n"
         "  AI能帮：画图/写方案/日照分析（5-10项）\n"
         "  AI做不到：客户沟通/审美判断/现场决策/签字负责\n\n"
         "\"会用AI的人会取代不会用AI的人\"",
         COLORS["accent"])

add_card(slide, 5.2, 1.5, 4.4, 1.8,
         "🚲 电动自行车类比",
         "\"AI是电动自行车，\n  不是自动驾驶汽车\"\n\n"
         "🔋 它放大你的力量\n"
         "    → 跑得更快、更远\n"
         "🧭 但你掌控方向\n"
         "    → 去哪、怎么走、何时停\n\n"
         "不是恐惧AI → 而是学会骑它",
         COLORS["green"])

add_card(slide, 0.4, 3.6, 9.2, 1.2,
         "📜 历史规律",
         "蒸汽机 → 淘汰手工纺织工 → 创造了铁路工程师、机械师  |  电力 → 淘汰煤气灯点灯人 → 创造了电工、电气工程师\n"
         "计算机 → 淘汰打字员 → 创造了程序员、UI设计师  |  互联网 → 淘汰部分零售店员 → 创造了电商运营、新媒体\n\n"
         "🔑 规律：技术淘汰\"可被标准化的重复性工作\"，创造\"需要人类判断力、创造力和情感智慧的工作\"",
         COLORS["secondary"])

add_card(slide, 0.4, 5.1, 9.2, 1.0,
         "📋 给大学生的四条建议",
         "1. 不要和AI比效率  |  2. 把AI当作\"初级助手\"  |  3. 培养AI做不了的能力（批判性思维/跨领域整合/审美判断/人际沟通）  |  4. 学好AI工具（不是可选项，是必选项）",
         COLORS["primary"])

add_teacher_tip(slide, "这个板块语气要诚恳但积极。核心：\"把焦虑转化为行动\"。如果学生质疑\"AI已经有创造力了\"，用Midjourney做例子回应。")
add_footer(slide, 10)

# ==========================================
# Slide 11: KB6 主流AI工具
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB6：主流AI工具介绍与选择", "国际六大工具 + 国内补充 + 决策树", "KB6", "10min")

# 六大工具卡片
tools = [
    ("ChatGPT", "OpenAI", "综合最强·生态完善\n$20/月·需特殊网络"),
    ("Claude", "Anthropic", "编程最强·长文分析\n$20/月·需特殊网络"),
    ("Gemini", "Google", "多模态最强·生态整合\n$20/月·需特殊网络"),
    ("⭐DeepSeek", "深度求索", "完全免费·中文极佳\n代码强·国内直连"),
    ("Kimi", "月之暗面", "200万字上下文\n论文利器·国内直连"),
    ("Perplexity", "Perplexity AI", "AI搜索引擎\n带引用来源·需特殊网络"),
]

for i, (name, company, desc) in enumerate(tools):
    left = 0.15 + i * 1.62
    top = 1.6
    is_featured = "DeepSeek" in name
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(1.52), Inches(1.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["yellow_bg"] if is_featured else COLORS["bg_light"]
    shape.line.color.rgb = COLORS["accent"] if is_featured else COLORS["secondary"]
    shape.line.width = Pt(2) if is_featured else Pt(1)
    add_textbox(slide, left + 0.1, top + 0.05, 1.32, 0.3, name,
                Pt(11), COLORS["primary"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.35, 1.32, 0.2, company,
                Pt(7), COLORS["text_light"], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.6, 1.32, 0.9, desc,
                Pt(8), COLORS["text_dark"])

# 国内补充 + 决策树
add_card(slide, 0.15, 3.55, 4.8, 2.0,
         "🇨🇳 其他国内工具",
         "• 通义千问（阿里）- 阿里生态·免费\n"
         "• 豆包（字节）- 语音交互好·免费\n"
         "• 文心一言（百度）- 百度生态·免费\n\n"
         "💰 零成本方案：\n"
         "DeepSeek + Kimi + 通义千问 + 豆包\n"
         "四个全免费·国内全直连",
         COLORS["green"])

add_card(slide, 5.2, 3.55, 4.6, 2.0,
         "🧭 选择决策树",
         "🖥 写代码 → Claude / DeepSeek\n"
         "📖 读论文 → Kimi（200万字）\n"
         "🔍 查最新信息 → Perplexity\n"
         "🖼 处理图片/视频 → Gemini\n"
         "🇨🇳 国内方便 → DeepSeek（首选）\n"
         "💰 全免费 → DeepSeek+Kimi+通义+豆包\n\n"
         "⚠️ 最重要：不要只用一个AI！\n"
         "组建你的\"AI顾问团\"，交叉验证",
         COLORS["secondary"])

add_teacher_tip(slide, "节奏要快！10分钟讲完。重点推荐DeepSeek作为学生首选工具（免费+国内直连+中文好）。")
add_footer(slide, 11)

# ==========================================
# Slide 12: KB7 实操环节 (上)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB7：动手实操 — 注册与体验AI工具", "本课最重要的环节！40分钟", "KB7", "40min")

# 实操步骤流程图
steps = [
    ("Step 1", "注册ChatGPT\n(15min)", "无法访问→跳过"),
    ("Step 2", "注册Claude\n(15min)", "无法访问→跳过"),
    ("Step 3", "注册Gemini\n(10min)", "无法访问→跳过"),
    ("Step 4", "体验Perplexity\n(10min)", "无法访问→跳过"),
    ("Step 5", "⭐ 注册国内AI\n(10min)", "DeepSeek·Kimi·通义"),
    ("Step 6", "⭐ 对比测试\n(最后完成)", "同一问题→多个AI"),
]

for i, (step, title, note) in enumerate(steps):
    left = 0.15 + i * 1.62
    top = 1.6
    is_key = "⭐" in title
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(1.52), Inches(1.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["orange_bg"] if is_key else COLORS["bg_light"]
    shape.line.color.rgb = COLORS["accent"] if is_key else COLORS["secondary"]
    shape.line.width = Pt(2) if is_key else Pt(1)
    add_textbox(slide, left + 0.1, top + 0.05, 1.32, 0.2, step,
                Pt(11), COLORS["accent"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.3, 1.32, 0.5, title,
                Pt(10), COLORS["text_dark"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.95, 1.32, 0.4, note,
                Pt(8), COLORS["text_mid"], alignment=PP_ALIGN.CENTER)
    # 箭头
    if i < 5:
        add_textbox(slide, left + 1.38, top + 0.55, 0.25, 0.3,
                    "→", Pt(12), COLORS["text_light"], alignment=PP_ALIGN.CENTER)

# 每种工具的关键信息
add_card(slide, 0.15, 3.4, 4.8, 1.6,
         "🔗 快速链接",
         "ChatGPT:   chatgpt.com\n"
         "Claude:    claude.ai\n"
         "Gemini:    gemini.google.com\n"
         "Perplexity: perplexity.ai\n"
         "⭐ DeepSeek: chat.deepseek.com\n"
         "Kimi:      kimi.moonshot.cn\n"
         "通义千问:   tongyi.aliyun.com\n"
         "豆包:      doubao.com",
         COLORS["primary"])

add_card(slide, 5.2, 3.4, 4.6, 1.6,
         "🎯 每个学生的最低目标",
         "✅ 至少成功注册1个AI工具\n"
         "   （推荐DeepSeek — 免费+国内直连）\n\n"
         "✅ 完成对比测试：\n"
         "   同一组问题 → 至少2个不同的AI\n"
         "   对比回答质量，记录差异\n\n"
         "⏱ 如果快：多注册几个工具\n"
         "⏱ 如果慢：只做DeepSeek + 对比测试",
         COLORS["accent"])

add_teacher_tip(slide, "实操期间教师不停在教室走动！看谁在皱眉/迷茫/卡住了。\"注册不了就跳过\"要用行动而非口头上落实。")
add_footer(slide, 12)

# ==========================================
# Slide 13: KB7 实操环节 (下) + 教师巡查
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB7：对比测试 + 教师巡查指南", "", "KB7", "40min")

add_card(slide, 0.4, 1.5, 4.4, 2.5,
         "🧪 统一测试Prompt（复制粘贴）",
         "Q1: 请用200字以内解释什么是\"结构力学\"。\n\n"
         "Q2: 一个工程设计师在日常工作中\n"
         "    如何使用AI提高效率？\n"
         "    请列出5个具体场景。\n\n"
         "Q3: 用Python写一个计算BMI的小程序\n"
         "    （输入身高m和体重kg，\n"
         "      输出BMI值和健康分级）。\n\n"
         "📋 发给至少2个不同的AI，记录对比",
         COLORS["primary"])

add_card(slide, 5.2, 1.5, 4.4, 2.5,
         "🔍 教师巡查五要点",
         "1. 有没有人眉头紧锁/发呆？\n"
         "   → 主动走过去：\"遇到什么问题？\"\n"
         "2. 有没有人卡在注册验证码？\n"
         "   → 立刻建议切换到DeepSeek\n"
         "3. 有没有人注册完了不知道干什么？\n"
         "   → 把测试Prompt发给他\n"
         "4. 有没有人AI回答看不懂？\n"
         "   → 教他追问：\"用更简单的语言解释\"\n"
         "5. 有没有人很快就完成了所有步骤？\n"
         "   → 让他帮助旁边的同学",
         COLORS["accent"])

add_card(slide, 0.4, 4.3, 9.2, 1.3,
         "⏱ 时间弹性方案",
         "如果时间充裕（50min）：让学生多问几个问题，深度体验AI工具的高级功能（文件上传、联网搜索等）\n"
         "如果时间紧张（30min）：砍掉Step 2-4（国外工具），全员只做 DeepSeek + 对比测试\n"
         "💡 口诀：\"境外工具全跳过，DeepSeek必须搞\"",
         COLORS["secondary"])

add_card(slide, 0.4, 5.8, 9.2, 0.7,
         "🆘 备用方案",
         "全班网络大面积故障 → 切换到\"全国内工具\"方案（DeepSeek + Kimi + 通义千问）  |  断网极端情况 → 让学生打开手机热点，或布置为课后作业",
         COLORS["red"])

add_teacher_tip(slide, "不要让任何一个学生在\"注册不了\"这件事上消耗超过2分钟。立刻给替代方案。")
add_footer(slide, 13)

# ==========================================
# Slide 14: KB8 课程总结与作业
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "KB8：课程总结与作业布置", "巩固核心认知 + 明确下周任务", "KB8", "5min")

add_card(slide, 0.4, 1.5, 4.4, 2.0,
         "📌 本课五句话核心回顾",
         "1️⃣ AI发展是一条70年的路，\n   不是一夜之间冒出来的\n\n"
         "2️⃣ ChatGPT = 技术突破 + 产品设计\n   + 时代机遇的完美结合\n\n"
         "3️⃣ 大模型能做的事情远超你想象\n   — 但需要你学会\"指挥\"它\n\n"
         "4️⃣ AI是电动自行车，不是自动驾驶\n   — 你掌握方向\n\n"
         "5️⃣ 不要只用一种AI\n   — 组建你的\"AI顾问团\"",
         COLORS["primary"])

add_card(slide, 5.2, 1.5, 4.4, 2.0,
         "📚 关键概念速查",
         "图灵测试：对话中分不清人还是机器\n"
         "Transformer：所有现代LLM的基础架构\n"
         "LLM：用海量文本训练的超大神经网络\n"
         "Prompt：你发给AI的指令/问题\n"
         "AI幻觉：AI一本正经地编造虚假信息\n"
         "Agent：能自主规划、使用工具的AI系统\n"
         "MCP：让AI安全调用外部工具的协议",
         COLORS["secondary"])

add_card(slide, 0.4, 3.8, 9.2, 2.0,
         "📝 课后作业（截止：下次上课前一天 22:00）",
         "作业1（必做）：AI工具对比笔记 — 3个问题 × 3个工具，≥400字，记录每个回答的优缺点\n"
         "作业2（必做）：传统方法 vs AI方法效率对比 — 选一个本专业真实问题，两种方式解决，记录时间/质量/信任度\n"
         "作业3（选做·推荐）：AI应用场景头脑风暴 — 列出5+个你未来工作中可以用AI提效的场景",
         COLORS["green"])

add_card(slide, 0.4, 6.0, 9.2, 0.6,
         "🔮 下节课预告",
         "Prompt Engineering（提示词工程）— RCTE框架 · 论文写作 · 代码编写 · PPT制作 · 方案设计 全部AI辅助",
         COLORS["accent"])

add_teacher_tip(slide, "5分钟收尾要快但不仓促。结束语要有仪式感：\"欢迎大家来到AI时代——这一次，你不是旁观者。\"")
add_footer(slide, 14)

# ==========================================
# Slide 15: 关键概念速查
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])
add_section_header(slide, "附录：本课关键概念速查表", "讲课过程中如有学生提问概念，可快速翻到此页", "参考", "")

concepts = [
    ("AI（人工智能）", "让计算机像人一样感知、思考、学习、决策的技术"),
    ("图灵测试", "如果机器能在对话中让人分不清它是人还是机器，就说明它有智能"),
    ("Transformer", "2017年提出的神经网络架构，所有现代大语言模型（GPT/Claude/Gemini）的基础"),
    ("大语言模型（LLM）", "用海量文本数据训练的超大型神经网络，能理解和生成人类语言"),
    ("ChatGPT", "OpenAI于2022年11月发布的对话式AI产品，引爆了全球AI浪潮"),
    ("Prompt（提示词）", "你发给AI的指令/问题（下周深入学习Prompt Engineering）"),
    ("多模态", "AI能同时处理文本、图片、音频、视频等多种类型的信息"),
    ("AI幻觉（Hallucination）", "AI一本正经地编造不真实的信息——大模型的主要局限性之一"),
    ("Agent（智能体）", "能自主规划、使用工具、完成复杂任务的AI系统（第7周重点学习）"),
    ("MCP", "Model Context Protocol，让AI能够安全调用外部工具的协议标准"),
    ("专用AI vs 通用AI", "专用AI只做一件事（如下棋），通用AI能处理多种任务（如ChatGPT）"),
    ("AI民主化", "2022年ChatGPT之后，AI从少数技术精英的工具变成所有人的工具"),
]

for i, (concept, desc) in enumerate(concepts):
    col = i // 6
    row = i % 6
    left = 0.2 + col * 5.0
    top = 1.5 + row * 0.85
    add_textbox(slide, left, top, 1.6, 0.25, concept,
                Pt(11), COLORS["primary"], bold=True)
    add_textbox(slide, left + 1.6, top, 3.2, 0.25, f"— {desc}",
                Pt(9), COLORS["text_dark"])

add_footer(slide, 15)

# ==========================================
# Slide 16: 结束页
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLORS["bg_white"])

# 底部色块
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.5), Inches(10), Inches(3.0)
)
shape.fill.solid()
shape.fill.fore_color.rgb = COLORS["primary"]
shape.line.fill.background()

add_textbox(slide, 1, 1.5, 8, 1.0,
            "72年前，图灵问：机器能思考吗？\n今天，你自己来回答这个问题。",
            Pt(24), COLORS["text_dark"], alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 2.8, 8, 0.8,
            "欢迎大家来到AI时代\n这一次，你不是旁观者。",
            Pt(28), COLORS["primary"], bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.0, 8, 0.5,
            "📅 下次课时间：[待定]  |  📍 [教室/线上链接]",
            Pt(16), COLORS["bg_white"], alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.6, 8, 0.5,
            "📋 作业截止：下次上课前一天 22:00  |  ❓ 课程群：[群名/群号]",
            Pt(14), RGBColor(0xC0, 0xD0, 0xE0), alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.2, 8, 0.5,
            "💬 课后有任何问题？用你刚注册的AI工具搜索答案！",
            Pt(14), RGBColor(0xC0, 0xD0, 0xE0), alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 7.0, 9.5, 0.3,
            "AI时代能力培养 · 第1课  |  V1.0  |  © 2026",
            Pt(8), RGBColor(0xA0, 0xB0, 0xC0), alignment=PP_ALIGN.CENTER)


# ==========================================
# 保存PPT
# ==========================================
output_dir = "d:/2026/AI-Learning Route/讲义/第1周/AI的发展与未来"
output_path = os.path.join(output_dir, "AI的发展与未来-教学课件.pptx")
prs.save(output_path)
print(f"✅ PPT已生成：{output_path}")
print(f"📊 共 {len(prs.slides)} 页幻灯片")
print(f"📁 文件大小：{os.path.getsize(output_path) / 1024:.1f} KB")
